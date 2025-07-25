#!/usr/bin/env python3
"""
PDF to PowerPoint Converter - Real conversion using PyMuPDF and python-pptx
Converts .pdf files to .pptx with OCR support for scanned PDFs
No AI services - local processing only
"""

import os
import sys
import json
import uuid
import time
import tempfile
from pathlib import Path

def log_message(message):
    """Log message with timestamp"""
    timestamp = time.strftime("%Y-%m-%d %H:%M:%S")
    print(f"[{timestamp}] {message}")

def detect_pdf_type(pdf_path):
    """Detect if PDF is text-based or scanned (image-based)"""
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(pdf_path)
        total_chars = 0
        total_pages = len(doc)
        
        # Sample first few pages to determine type
        sample_pages = min(3, total_pages)
        
        for page_num in range(sample_pages):
            page = doc[page_num]
            text = page.get_text()
            total_chars += len(text.strip())
        
        doc.close()
        
        # If less than 50 characters per page on average, likely scanned
        avg_chars_per_page = total_chars / sample_pages if sample_pages > 0 else 0
        is_scanned = avg_chars_per_page < 50
        
        log_message(f"PDF analysis: {total_pages} pages, {avg_chars_per_page:.1f} chars/page avg, {'scanned' if is_scanned else 'text-based'}")
        
        return is_scanned, total_pages, avg_chars_per_page
        
    except Exception as e:
        log_message(f"PDF type detection failed: {str(e)}")
        return False, 0, 0

def extract_text_from_pdf(pdf_path):
    """Extract text from text-based PDF using PyMuPDF"""
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(pdf_path)
        pages_text = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            # Split text into logical sections
            paragraphs = []
            for line in text.split('\n'):
                line = line.strip()
                if line:
                    paragraphs.append(line)
            
            pages_text.append({
                'page_number': page_num + 1,
                'text': text,
                'paragraphs': paragraphs
            })
        
        doc.close()
        log_message(f"Extracted text from {len(pages_text)} pages")
        return pages_text
        
    except Exception as e:
        log_message(f"Text extraction failed: {str(e)}")
        return []

def extract_text_ocr(pdf_path):
    """Extract text from scanned PDF using OCR"""
    try:
        from pdf2image import convert_from_path
        import pytesseract
        from PIL import Image
        
        log_message("Converting PDF pages to images for OCR...")
        images = convert_from_path(pdf_path)
        pages_text = []
        
        for i, image in enumerate(images):
            log_message(f"Processing page {i+1} with OCR...")
            
            # Perform OCR on the image
            text = pytesseract.image_to_string(image, lang='eng')
            
            # Process OCR text into paragraphs
            paragraphs = []
            for line in text.split('\n'):
                line = line.strip()
                if line and len(line) > 3:  # Filter out very short lines
                    paragraphs.append(line)
            
            pages_text.append({
                'page_number': i + 1,
                'text': text,
                'paragraphs': paragraphs
            })
        
        log_message(f"OCR completed for {len(pages_text)} pages")
        return pages_text
        
    except ImportError:
        log_message("OCR libraries not available (pdf2image, pytesseract)")
        return []
    except Exception as e:
        log_message(f"OCR extraction failed: {str(e)}")
        return []

def create_powerpoint(pages_text, output_path):
    """Create PowerPoint presentation from extracted text"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "PDF to PowerPoint Conversion"
        subtitle.text = f"Generated from PDF • {len(pages_text)} pages"
        
        # Process each page
        for page_data in pages_text:
            # Add a new slide for each page
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            title_shape = slide.shapes.title
            title_shape.text = f"Page {page_data['page_number']}"
            
            # Add content
            if page_data['paragraphs']:
                # Use content placeholder if available
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    # Add paragraphs as bullet points
                    for i, paragraph_text in enumerate(page_data['paragraphs'][:10]):  # Limit to 10 items
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = paragraph_text
                        p.level = 0
                        
                        # Formatting
                        font = p.font
                        font.size = Pt(14)
                        font.name = 'Arial'
                else:
                    # Add text box if no content placeholder
                    left = Inches(1)
                    top = Inches(1.5)
                    width = Inches(8)
                    height = Inches(5)
                    
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    
                    for paragraph_text in page_data['paragraphs'][:15]:  # Limit to 15 items
                        p = text_frame.add_paragraph()
                        p.text = f"• {paragraph_text}"
                        font = p.font
                        font.size = Pt(12)
                        font.name = 'Arial'
            else:
                # No content found
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = "No text content extracted from this page."
        
        # Save presentation
        prs.save(output_path)
        
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            log_message(f"PowerPoint created successfully: {output_path}")
            return True, f"Created presentation with {len(pages_text)} slides"
        else:
            return False, "PowerPoint file was not created or is empty"
            
    except ImportError:
        return False, "python-pptx library not available"
    except Exception as e:
        return False, f"PowerPoint creation failed: {str(e)}"

def main():
    """Main conversion function"""
    if len(sys.argv) != 3:
        print(json.dumps({
            "success": False,
            "error": "Usage: python pdf_to_pptx_converter.py <input_pdf> <output_dir>"
        }))
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_dir = sys.argv[2]
    
    try:
        # Validate input file
        if not os.path.exists(input_path):
            print(json.dumps({
                "success": False,
                "error": "Input file does not exist"
            }))
            sys.exit(1)
        
        # Check file extension
        if not input_path.lower().endswith('.pdf'):
            print(json.dumps({
                "success": False,
                "error": "Only .pdf files are supported"
            }))
            sys.exit(1)
        
        # Check file size (50MB limit)
        file_size = os.path.getsize(input_path)
        if file_size > 50 * 1024 * 1024:
            print(json.dumps({
                "success": False,
                "error": "File size exceeds 50MB limit"
            }))
            sys.exit(1)
        
        log_message(f"Processing PDF: {input_path} ({file_size} bytes)")
        
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        # Detect PDF type
        is_scanned, total_pages, avg_chars = detect_pdf_type(input_path)
        
        if total_pages == 0:
            print(json.dumps({
                "success": False,
                "error": "Could not read PDF file or file is empty"
            }))
            sys.exit(1)
        
        # Extract text based on PDF type
        if is_scanned:
            log_message("PDF appears to be scanned - using OCR")
            pages_text = extract_text_ocr(input_path)
            extraction_method = "OCR"
        else:
            log_message("PDF appears to be text-based - extracting text directly")
            pages_text = extract_text_from_pdf(input_path)
            extraction_method = "Direct text extraction"
        
        if not pages_text:
            print(json.dumps({
                "success": False,
                "error": "Failed to extract text from PDF"
            }))
            sys.exit(1)
        
        # Generate output path
        input_stem = Path(input_path).stem
        output_path = os.path.join(output_dir, f"{input_stem}.pptx")
        
        # Create PowerPoint presentation
        success, message = create_powerpoint(pages_text, output_path)
        
        if success:
            output_size = os.path.getsize(output_path)
            
            print(json.dumps({
                "success": True,
                "output_path": output_path,
                "file_size": output_size,
                "total_pages": total_pages,
                "extraction_method": extraction_method,
                "message": message,
                "is_scanned": is_scanned
            }))
        else:
            print(json.dumps({
                "success": False,
                "error": message
            }))
            sys.exit(1)
    
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": f"Unexpected error: {str(e)}"
        }))
        sys.exit(1)

if __name__ == "__main__":
    main()